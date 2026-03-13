"""
PPT 导出 Skill（专业版）
=======================
使用 python-pptx 生成专业品质的 PowerPoint 演示文稿。
特性：
- 自定义专业配色方案（深蓝+科技蓝主色调）
- 多种幻灯片布局（封面、章节页、内容页、总结页）
- Markdown 智能解析为结构化幻灯片
- 自动分页（超长内容拆分到多张幻灯片）
- 统一的字体与排版规范
"""

import os
import re
import logging
from typing import List, Tuple

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

from skills.export_utils import get_export_dir, build_combined_markdown

logger = logging.getLogger("agent_skills.export_pptx")

# ── 专业配色方案 ──────────────────────────────────────
COLOR_PRIMARY = RGBColor(0x1A, 0x56, 0xDB)      # 主色：科技蓝
COLOR_PRIMARY_DARK = RGBColor(0x0F, 0x3A, 0x8C)  # 深蓝
COLOR_ACCENT = RGBColor(0x2E, 0xCC, 0x71)        # 强调色：翠绿
COLOR_TEXT_DARK = RGBColor(0x1F, 0x2D, 0x3D)     # 正文深色
COLOR_TEXT_LIGHT = RGBColor(0xFF, 0xFF, 0xFF)     # 白色文字
COLOR_TEXT_GRAY = RGBColor(0x60, 0x6F, 0x7B)     # 灰色辅助文字
COLOR_BG_LIGHT = RGBColor(0xF0, 0xF4, 0xF8)      # 浅灰背景
COLOR_BORDER = RGBColor(0xD4, 0xDE, 0xE8)         # 边框色

# ── 幻灯片尺寸与边距 ─────────────────────────────────
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
MARGIN_LEFT = Inches(0.8)
MARGIN_RIGHT = Inches(0.8)
MARGIN_TOP = Inches(0.6)
CONTENT_WIDTH = Inches(11.733)

# ── 每张内容页最多显示的文本行数 ─────────────────────
MAX_BULLETS_PER_SLIDE = 7
MAX_CHARS_PER_BULLET = 120


def export_to_pptx(
    session_id: str,
    project_name: str,
    contents: List[Tuple[str, str]],
    filename_suffix: str = "_需求与详设",
) -> Tuple[bool, str, str]:
    """
    将多段 (标题, 内容) 转为专业品质的 PPT 演示文稿。

    Returns:
        (success, file_path, error_message)
    """
    export_dir = get_export_dir(session_id)
    filename = f"{project_name or 'document'}{filename_suffix}.pptx"
    filename = _sanitize_filename(filename)
    file_path = os.path.join(export_dir, filename)

    try:
        prs = Presentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT

        _add_cover_slide(prs, project_name)

        for title, body in contents:
            if not body or not body.strip():
                continue
            _add_section_slide(prs, title)
            sections = _parse_markdown_to_sections(body)
            for section in sections:
                _add_content_slides(prs, section["heading"], section["items"])

        _add_ending_slide(prs, project_name)

        prs.save(file_path)
        logger.info("PPTX 已保存: %s", file_path)
        return True, file_path, ""
    except Exception as exc:
        logger.exception("PPTX 导出失败: %s", exc)
        return False, "", str(exc)


# ── 封面页 ──────────────────────────────────────────

def _add_cover_slide(prs: Presentation, project_name: str) -> None:
    """添加专业封面页：深蓝背景 + 白色标题。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_slide_background(slide, COLOR_PRIMARY_DARK)

    # 左侧装饰色条
    bar = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(0.3), SLIDE_HEIGHT
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_ACCENT
    bar.line.fill.background()

    # 主标题
    title_box = slide.shapes.add_textbox(
        Inches(1.2), Inches(2.2), Inches(10), Inches(1.5)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = project_name or "需求与详设文档"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_LIGHT
    p.alignment = PP_ALIGN.LEFT

    # 副标题
    sub_box = slide.shapes.add_textbox(
        Inches(1.2), Inches(4.0), Inches(10), Inches(0.8)
    )
    sub_tf = sub_box.text_frame
    sub_p = sub_tf.paragraphs[0]
    sub_p.text = "AI 需求与详设生成器 · Doc-genius"
    sub_p.font.size = Pt(18)
    sub_p.font.color.rgb = RGBColor(0xA0, 0xC4, 0xF0)
    sub_p.alignment = PP_ALIGN.LEFT

    # 底部装饰线
    line = slide.shapes.add_shape(
        1, Inches(1.2), Inches(3.7), Inches(3), Pt(3)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_ACCENT
    line.line.fill.background()


# ── 章节分隔页 ──────────────────────────────────────

def _add_section_slide(prs: Presentation, section_title: str) -> None:
    """添加章节分隔页：蓝色背景 + 居中标题。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_slide_background(slide, COLOR_PRIMARY)

    # 章节标题
    title_box = slide.shapes.add_textbox(
        MARGIN_LEFT, Inches(2.5), CONTENT_WIDTH, Inches(1.5)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_LIGHT
    p.alignment = PP_ALIGN.CENTER

    # 底部装饰线
    line_left = (SLIDE_WIDTH - Inches(4)) // 2
    line = slide.shapes.add_shape(
        1, line_left, Inches(4.2), Inches(4), Pt(3)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_ACCENT
    line.line.fill.background()


# ── 内容页 ──────────────────────────────────────────

def _add_content_slides(
    prs: Presentation,
    heading: str,
    items: List[str],
) -> None:
    """
    添加内容幻灯片（白色背景 + 蓝色顶栏）。
    超长内容自动拆分到多张幻灯片。
    """
    if not items:
        items = ["（暂无详细内容）"]

    chunks = _chunk_items(items, MAX_BULLETS_PER_SLIDE)
    total_pages = len(chunks)

    for page_idx, chunk in enumerate(chunks):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 顶部蓝色横条
        top_bar = slide.shapes.add_shape(
            1, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.9)
        )
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = COLOR_PRIMARY_DARK
        top_bar.line.fill.background()

        # 页面标题（顶栏内）
        page_label = f" ({page_idx + 1}/{total_pages})" if total_pages > 1 else ""
        title_box = slide.shapes.add_textbox(
            MARGIN_LEFT, Inches(0.15), CONTENT_WIDTH, Inches(0.6)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = heading + page_label
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_LIGHT
        p.alignment = PP_ALIGN.LEFT

        # 内容区域
        content_box = slide.shapes.add_textbox(
            MARGIN_LEFT, Inches(1.2), CONTENT_WIDTH, Inches(5.8)
        )
        content_box.text_frame.word_wrap = True

        for i, item in enumerate(chunk):
            if i == 0:
                p = content_box.text_frame.paragraphs[0]
            else:
                p = content_box.text_frame.add_paragraph()

            # 检测是否为子项（以 - 或 · 开头的二级缩进）
            is_sub = item.startswith("  ") or item.startswith("\t")
            clean_text = item.strip().lstrip("-·•").strip()

            if is_sub:
                p.text = f"    · {clean_text}"
                p.font.size = Pt(14)
                p.font.color.rgb = COLOR_TEXT_GRAY
            else:
                p.text = f"▸ {clean_text}"
                p.font.size = Pt(16)
                p.font.color.rgb = COLOR_TEXT_DARK

            p.space_after = Pt(6)
            p.alignment = PP_ALIGN.LEFT

        # 左侧蓝色装饰线
        accent_line = slide.shapes.add_shape(
            1, Inches(0.4), Inches(1.2), Pt(4), Inches(5.5)
        )
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = COLOR_PRIMARY
        accent_line.line.fill.background()


# ── 总结/结束页 ─────────────────────────────────────

def _add_ending_slide(prs: Presentation, project_name: str) -> None:
    """添加结束页：深色背景 + 感谢语。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_slide_background(slide, COLOR_PRIMARY_DARK)

    title_box = slide.shapes.add_textbox(
        MARGIN_LEFT, Inches(2.4), CONTENT_WIDTH, Inches(1.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_LIGHT
    p.alignment = PP_ALIGN.CENTER

    sub_box = slide.shapes.add_textbox(
        MARGIN_LEFT, Inches(4.2), CONTENT_WIDTH, Inches(0.8)
    )
    sub_tf = sub_box.text_frame
    sub_p = sub_tf.paragraphs[0]
    sub_p.text = f"{project_name} — 由 Doc-genius AI 生成"
    sub_p.font.size = Pt(16)
    sub_p.font.color.rgb = RGBColor(0xA0, 0xC4, 0xF0)
    sub_p.alignment = PP_ALIGN.CENTER


# ── Markdown 解析 ───────────────────────────────────

def _parse_markdown_to_sections(md_text: str) -> List[dict]:
    """
    将 Markdown 文本解析为结构化的幻灯片数据。

    返回:
        [{"heading": "标题", "items": ["要点1", "要点2", ...]}, ...]
    """
    sections: List[dict] = []
    current_heading = ""
    current_items: List[str] = []

    for line in md_text.split("\n"):
        stripped = line.strip()
        if not stripped or stripped.startswith("---"):
            continue

        # 标题行
        heading_match = re.match(r"^(#{1,4})\s+(.+)$", stripped)
        if heading_match:
            if current_heading or current_items:
                sections.append({
                    "heading": current_heading or "概述",
                    "items": current_items,
                })
            current_heading = heading_match.group(2).strip()
            current_items = []
            continue

        # 列表项（有序/无序）
        list_match = re.match(r"^[-*]\s+(.+)$", stripped)
        ordered_match = re.match(r"^\d+\.\s+(.+)$", stripped)
        if list_match:
            current_items.append(list_match.group(1).strip())
        elif ordered_match:
            current_items.append(ordered_match.group(1).strip())
        elif stripped.startswith("|") and "|" in stripped[1:]:
            # 表格行 → 转为普通文本
            cells = [c.strip() for c in stripped.strip("|").split("|") if c.strip() and not re.match(r"^[-:]+$", c.strip())]
            if cells:
                current_items.append(" | ".join(cells))
        else:
            # 普通段落 → 作为内容项
            text = _strip_md_formatting(stripped)
            if text and len(text) > 2:
                current_items.append(text)

    if current_heading or current_items:
        sections.append({
            "heading": current_heading or "概述",
            "items": current_items,
        })

    if not sections:
        sections.append({"heading": "内容", "items": ["（暂无详细内容）"]})

    return sections


def _strip_md_formatting(text: str) -> str:
    """移除 Markdown 行内格式（加粗、斜体、行内代码等）。"""
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"\*(.+?)\*", r"\1", text)
    text = re.sub(r"`(.+?)`", r"\1", text)
    text = re.sub(r"\[(.+?)\]\(.+?\)", r"\1", text)
    return text.strip()


# ── 工具函数 ────────────────────────────────────────

def _chunk_items(items: List[str], max_per_page: int) -> List[List[str]]:
    """将列表按最大条数拆分为多页。"""
    if not items:
        return [[]]
    chunks = []
    for i in range(0, len(items), max_per_page):
        chunks.append(items[i:i + max_per_page])
    return chunks


def _fill_slide_background(slide, color: RGBColor) -> None:
    """设置幻灯片纯色背景。"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _sanitize_filename(name: str) -> str:
    """移除文件名中的非法字符。"""
    for c in ["/", "\\", ":", "*", "?", '"', "<", ">", "|"]:
        name = name.replace(c, "_")
    return name[:200] if len(name) > 200 else name
